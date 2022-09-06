#! /usr/bin/python3

import copy
from pptx import Presentation
from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart

# Helper function that returns the layout type of a 
# presentation
def get_slide_layout(prs):
	num_layouts = [len(layout.placeholders) for layout in prs.slide_layouts]
	min_size = min(num_layouts)
	index = num_layouts.index(min_size)
	return prs.slide_layouts[index]

# Helper function that copies over slide data from a
# source to a destination presentation using slide 
# indexing. Slides are inserted at the end of the
# destination presentation
def copy_slide(slide, dst, layout):
	# Retrieve slide from source PPT and create empty
	# slide in destination PPT 
	in_slide = slide
	out_slide = dst.slides.add_slide(layout)

	# Copy over slide elements 
	for shape in in_slide.shapes:
		element = copy.deepcopy(shape.element)
		out_slide.shapes._spTree.insert_element_before(element, 'p:estLst')

	# Update element positions
	for key, value in in_slide.part.rels.items():
		# Prevent the copying of a notesSlide relation 
		if not 'notesSlide' in value.reltype:
			target = value._target
			# Handle charts by creating a new embedded chart
			if "chart" in value.reltype:
				partname = target.package.next_partname(ChartPart.partname_template)
				xlsx_blob = target.chart_workbook.xlsx_part.blob
				target = ChartPart(partname, target.content_type, copy.deepcopy(target._element), package=target.package)

				target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(xlsx_blob, target.package)
			# Update element relationships
			out_slide.part.rels.add_relationship(value.reltype, target, value.rId)

def main():
	# Get user input
	# num_merge = int(input("Please enter the number of PPTs you want to merge: "))
	# outfile = input("Please enter a name for the merged file: ")
	# fpath_1 = input("Please enter the path to the first file: ")

	# Test filepaths
	fpath_1 = "./test3.pptx"
	fpath_2 = "./test4.pptx"

	# Open files for reading
	prs_main = Presentation(fpath_1)
	prs_next = Presentation(fpath_2)

	# Get layout of source slides
	slide_layout = get_slide_layout(prs_main)

	# Loop that iterates through slides to copy
	for slide in prs_next.slides:
		copy_slide(slide, prs_main, slide_layout)

	# Save output
	prs_main.save('merged.pptx')

	# for _ in range(1, num_merge-1):
	# 	fpath_2 = input("Please enter the path to the next file: ")
	# 	print(fpath2)

if __name__ == '__main__':
	main()